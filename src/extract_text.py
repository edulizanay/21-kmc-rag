# ABOUTME: Extracts text from binary document formats (docx, xlsx, pptx, pdf, odt).
# ABOUTME: Used by triage subagents to read file contents they can't access via Read tool.

import sys
from pathlib import Path


def extract_text(filepath: str, max_chars: int = 5000) -> str:
    """Extract text from a file, truncated to max_chars."""
    path = Path(filepath)
    ext = path.suffix.lower()

    try:
        if ext == ".docx":
            return _extract_docx(path, max_chars)
        elif ext == ".xlsx":
            return _extract_xlsx(path, max_chars)
        elif ext == ".pptx":
            return _extract_pptx(path, max_chars)
        elif ext == ".pdf":
            return _extract_pdf(path, max_chars)
        elif ext == ".csv":
            return _extract_csv(path, max_chars)
        elif ext in (".md", ".txt", ".yaml", ".json", ".odt"):
            return path.read_text(encoding="utf-8", errors="replace")[:max_chars]
        else:
            return f"[Unsupported format: {ext}]"
    except Exception as e:
        return f"[Error extracting {ext}: {e}]"


def _extract_docx(path: Path, max_chars: int) -> str:
    from docx import Document

    doc = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return text[:max_chars]


def _extract_xlsx(path: Path, max_chars: int) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[Sheet: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
            if len("\n".join(parts)) > max_chars:
                break
        if len("\n".join(parts)) > max_chars:
            break
    wb.close()
    return "\n".join(parts)[:max_chars]


def _extract_pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
        if len("\n".join(parts)) > max_chars:
            break
    return "\n".join(parts)[:max_chars]


def _extract_pdf(path: Path, max_chars: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        parts.append(text)
        if len("\n".join(parts)) > max_chars:
            break
    return "\n".join(parts)[:max_chars]


def _extract_csv(path: Path, max_chars: int) -> str:
    return path.read_text(encoding="utf-8", errors="replace")[:max_chars]


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_text.py <filepath> [max_chars]")
        sys.exit(1)
    filepath = sys.argv[1]
    max_chars = int(sys.argv[2]) if len(sys.argv) > 2 else 5000
    print(extract_text(filepath, max_chars))
